"""Gera a apresentacao do projeto (PPTX) de forma reproduzivel.

Uso:
    pip install python-pptx
    python docs/apresentacao/gerar_apresentacao.py

Saida: docs/apresentacao/Localizacao-Indoor-Teto.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Paleta
AZUL = RGBColor(0x1F, 0x29, 0x37)
AZUL_CLARO = RGBColor(0x25, 0x63, 0xEB)
VERDE = RGBColor(0x16, 0xA3, 0x4A)
CINZA = RGBColor(0x6B, 0x72, 0x80)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ESCURO = RGBColor(0x11, 0x18, 0x27)

LARGURA = Inches(13.333)
ALTURA = Inches(7.5)


def _fundo(slide, cor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = cor


def _caixa(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def slide_capa(prs, titulo, subtitulo, rodape):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(s, AZUL)
    tf = _caixa(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.5))
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p2 = tf.add_paragraph()
    p2.text = subtitulo
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
    tf3 = _caixa(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.8))
    pr = tf3.paragraphs[0]
    pr.text = rodape
    pr.font.size = Pt(14)
    pr.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)


def slide_topicos(prs, titulo, topicos, cor_titulo=AZUL_CLARO):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(s, BRANCO)
    # Barra do titulo
    tf = _caixa(s, Inches(0.7), Inches(0.5), Inches(12), Inches(1.0))
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = cor_titulo
    # Corpo
    corpo = _caixa(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.4))
    primeiro = True
    for item in topicos:
        nivel = 0
        texto = item
        if isinstance(item, tuple):
            texto, nivel = item
        p = corpo.paragraphs[0] if primeiro else corpo.add_paragraph()
        primeiro = False
        marcador = "•" if nivel == 0 else "–"
        p.text = f"{'   ' * nivel}{marcador} {texto}"
        p.font.size = Pt(22 - nivel * 3)
        p.font.color.rgb = ESCURO if nivel == 0 else CINZA
        p.space_after = Pt(8)


def slide_secao(prs, titulo, cor=VERDE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(s, cor)
    tf = _caixa(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.5))
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.LEFT


def slide_tabela(prs, titulo, cabecalho, linhas):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(s, BRANCO)
    tf = _caixa(s, Inches(0.7), Inches(0.5), Inches(12), Inches(1.0))
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = AZUL_CLARO

    n_lin = len(linhas) + 1
    n_col = len(cabecalho)
    tabela = s.shapes.add_table(
        n_lin, n_col, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4 * n_lin)
    ).table
    for j, texto in enumerate(cabecalho):
        cel = tabela.cell(0, j)
        cel.text = texto
        cel.fill.solid()
        cel.fill.fore_color.rgb = AZUL
        par = cel.text_frame.paragraphs[0]
        par.font.bold = True
        par.font.size = Pt(16)
        par.font.color.rgb = BRANCO
    for i, linha in enumerate(linhas, start=1):
        for j, texto in enumerate(linha):
            cel = tabela.cell(i, j)
            cel.text = texto
            par = cel.text_frame.paragraphs[0]
            par.font.size = Pt(14)
            par.font.color.rgb = ESCURO


def construir() -> Path:
    prs = Presentation()
    prs.slide_width = LARGURA
    prs.slide_height = ALTURA

    slide_capa(
        prs,
        "Localização Indoor Baseada em Imagens do Teto",
        "Estimar onde você está, dentro de um prédio, a partir de uma foto do teto + GPS",
        "Projeto de pesquisa · Munif Gebara · 2026",
    )

    slide_topicos(
        prs,
        "O problema",
        [
            "GPS funciona mal dentro de prédios (paredes, lajes, múltiplos andares).",
            "Localizar pessoas/ativos em shoppings, aeroportos e hospitais é difícil.",
            "Soluções com Wi-Fi/Bluetooth exigem infraestrutura e calibração constante.",
            ("Queremos algo que use só o que o celular já tem: câmera + GPS.", 1),
        ],
    )

    slide_topicos(
        prs,
        "A hipótese: por que o teto?",
        [
            "O teto é uma referência visual estável:",
            ("Quase não tem pessoas passando na frente.", 1),
            ("Muda pouco com o tempo.", 1),
            ("Iluminação mais constante que no nível do chão.", 1),
            "Logo, é ideal para reconhecimento visual de lugar (Visual Place Recognition).",
            "GPS impreciso é suficiente como filtro grosseiro (reduz o espaço de busca).",
        ],
    )

    slide_topicos(
        prs,
        "Como funciona (visão geral)",
        [
            "Fase 1 — Mapeamento:",
            ("Coletor tira fotos do teto; a cada foto captura o GPS.", 1),
            ("Usuário toca na planta baixa para marcar onde a foto foi tirada (x, y).", 1),
            ("Cada foto vira um vetor (embedding) guardado no banco.", 1),
            "Fase 2 — Consulta:",
            ("Nova foto + GPS chegam na API.", 1),
            ("A API filtra por GPS, acha as fotos mais parecidas e devolve (x, y) + confiança.", 1),
        ],
    )

    slide_topicos(
        prs,
        "Pipeline de localização (o núcleo)",
        [
            "1. Filtro geográfico por raio de GPS  →  milhares viram dezenas (earthdistance).",
            "2. Busca por similaridade visual  →  top-N candidatos (pgvector, embeddings).",
            "3. Verificação geométrica  →  features locais + RANSAC confirmam o casamento.",
            "4. Interpolação ponderada  →  estima (x, y) na planta a partir dos melhores.",
            "5. Resposta  →  { x, y, confiança, candidatos }.",
        ],
    )

    slide_tabela(
        prs,
        "Decisões de arquitetura",
        ["Tema", "Decisão"],
        [
            ["Backend", "Python + FastAPI"],
            ["Match visual", "Híbrido: embeddings + verificação geométrica"],
            ["Banco", "PostgreSQL + pgvector + earthdistance/cube"],
            ["Saída da API", "(x, y) interpolado na planta + confiança"],
            ["Planta baixa", "SVG canônico (DWG convertido no upload)"],
            ["Frontend", "Vue 3 + Vite (PWA): admin, coletor, teste"],
            ["Deploy", "Kubernetes (k3s) no gimli + Cloudflare Tunnel"],
        ],
    )

    slide_topicos(
        prs,
        "Modelo de dados",
        [
            "local — o ambiente mapeado (shopping, aeroporto...).",
            "planta — planta baixa em SVG vinculada a um local.",
            "foto — foto do teto:",
            ("latitude, longitude, gps_precisao (GPS bruto)", 1),
            ("plan_x, plan_y (posição marcada na planta)", 1),
            ("embedding vector(512) (pgvector)", 1),
            ("tipo: 'mapeamento' ou 'consulta'", 1),
            "Índices: GiST em ll_to_earth(lat,lon) e HNSW no embedding.",
        ],
    )

    slide_topicos(
        prs,
        "Os três aplicativos (PWA)",
        [
            "Admin — cadastra locais e faz upload da planta (SVG).",
            "Coletor — tira foto do teto, captura GPS e marca a posição na planta.",
            "Teste — envia foto + GPS e mostra a posição estimada.",
            "São PWAs em Vue 3: instaláveis, acessam câmera e GPS pelo navegador.",
            ("Câmera/GPS são Web APIs — sem app nativo.", 1),
            "Código comum (câmera, GPS, HTTP, planta) compartilhado entre os três.",
        ],
    )

    slide_topicos(
        prs,
        "Infraestrutura e deploy",
        [
            "Servidor gimli (192.168.0.99): Kubernetes k3s single-node.",
            "Postgres compartilhado (pgvector/pgvector:pg17) — banco 'teto' dedicado.",
            "Ingress Traefik; storage local-path.",
            "Imagem construída no nó e importada no containerd (sem registry externo).",
            "Exposição externa via Cloudflare Tunnel.",
            "CI/CD com GitHub Actions self-hosted runner.",
            "Fluxo: issue → branch → Pull Request → merge na master.",
        ],
    )

    slide_topicos(
        prs,
        "Status atual — M0 concluído ✓",
        [
            "Repositório, estrutura e documentação (em português) prontos.",
            "Backend FastAPI no ar no gimli: /saude → banco 'ok'.",
            "Banco 'teto' criado com extensões vector, cube, earthdistance.",
            "Tabelas local, planta, foto + índices (HNSW, GiST) via migração Alembic.",
            "Apps Vue (admin, coletor, teste) com estrutura e integração câmera/GPS.",
            ("Acessível em teto.local (Traefik); Cloudflare Tunnel a configurar.", 1),
        ],
        cor_titulo=VERDE,
    )

    slide_tabela(
        prs,
        "Roadmap",
        ["Marco", "Entrega"],
        [
            ["M0 ✓", "Fundação: infra, API de saúde, banco e deploy no gimli"],
            ["M1", "Cadastro: API + Admin (criar local, upload da planta SVG)"],
            ["M2", "Coleta: Coletor grava foto + GPS + (x,y) + embedding"],
            ["M3", "Localização: endpoint de consulta + Teste mostra o ponto"],
            ["M4", "Avaliação: métricas de erro (metros) e ajuste do modelo"],
        ],
    )

    slide_secao(prs, "Obrigado!  Próximo passo: M1 — Cadastro de locais e plantas", cor=AZUL_CLARO)

    saida = Path(__file__).parent / "Localizacao-Indoor-Teto.pptx"
    prs.save(saida)
    return saida


if __name__ == "__main__":
    caminho = construir()
    print(f"Apresentacao gerada: {caminho}")
